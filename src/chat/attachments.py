import asyncio
import base64
import csv
import gzip
import io
import json
import mimetypes
import re
import tarfile
import zipfile
from dataclasses import dataclass
from email import policy
from email.parser import BytesParser
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

from defusedxml import ElementTree
from docx import Document
from openpyxl import load_workbook
from PIL import Image, ImageOps
from pillow_heif import register_heif_opener
from pptx import Presentation
from pypdf import PdfReader

register_heif_opener()

_TEXT_EXTENSIONS = {
    ".asm",
    ".bat",
    ".bash",
    ".c",
    ".cfg",
    ".conf",
    ".cpp",
    ".cs",
    ".css",
    ".csv",
    ".go",
    ".graphql",
    ".h",
    ".hpp",
    ".htm",
    ".html",
    ".ini",
    ".ipynb",
    ".java",
    ".js",
    ".json",
    ".jsx",
    ".log",
    ".lua",
    ".md",
    ".mjs",
    ".pl",
    ".php",
    ".properties",
    ".py",
    ".rb",
    ".rs",
    ".rst",
    ".sh",
    ".sql",
    ".svg",
    ".swift",
    ".tex",
    ".toml",
    ".ts",
    ".tsx",
    ".txt",
    ".xml",
    ".yaml",
    ".yml",
}
_ARCHIVE_EXTENSIONS = {".gz", ".tar", ".tgz", ".zip"}
_IMAGE_EXTENSIONS = {
    ".avif",
    ".bmp",
    ".gif",
    ".heic",
    ".heif",
    ".jpeg",
    ".jpg",
    ".png",
    ".tif",
    ".tiff",
    ".webp",
}
_OPEN_DOCUMENT_EXTENSIONS = {".odt", ".ods", ".odp"}
_MAX_IMAGE_PIXELS = 40_000_000
_MAX_IMAGE_EDGE = 4096
_MAX_ARCHIVE_MEMBERS = 50
_MAX_ARCHIVE_EXPANDED_BYTES = 20 * 1024 * 1024


class AttachmentError(ValueError):
    pass


@dataclass(frozen=True, slots=True)
class ProcessedAttachments:
    text: str
    images: list[str]
    names: list[str]


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: list[str] = []

    def handle_data(self, data: str) -> None:
        stripped = data.strip()
        if stripped:
            self.parts.append(stripped)


async def process_attachments(
    elements: list[Any],
    *,
    max_files: int,
    max_file_bytes: int,
    max_extracted_chars: int,
) -> ProcessedAttachments:
    if len(elements) > max_files:
        raise AttachmentError(f"Upload at most {max_files} files at a time.")

    sections: list[str] = []
    images: list[str] = []
    names: list[str] = []
    remaining = max_extracted_chars
    for element in elements:
        name = _element_name(element)
        payload = await _read_element(element, max_file_bytes)
        extension = Path(name).suffix.casefold()
        mime = str(getattr(element, "mime", "") or "")
        if extension in _IMAGE_EXTENSIONS or mime.startswith("image/"):
            encoded = await asyncio.to_thread(_validate_and_encode_image, payload)
            images.append(encoded)
            names.append(name)
            continue

        extracted = await asyncio.to_thread(_extract_text, name, mime, payload)
        if not extracted.strip():
            raise AttachmentError(f"No readable text was found in {name}.")
        excerpt = extracted[:remaining]
        sections.append(f"--- Attachment: {name} ---\n{excerpt}")
        names.append(name)
        remaining -= len(excerpt)
        if remaining <= 0:
            break

    return ProcessedAttachments(
        text="\n\n".join(sections),
        images=images,
        names=names,
    )


def _element_name(element: Any) -> str:
    name = str(getattr(element, "name", "") or "").strip()
    return Path(name).name or "attachment"


async def _read_element(element: Any, max_file_bytes: int) -> bytes:
    content = getattr(element, "content", None)
    if isinstance(content, bytes):
        payload = content
    else:
        path_value = getattr(element, "path", None)
        if not path_value:
            raise AttachmentError(f"Cannot access {_element_name(element)}.")
        path = Path(str(path_value))
        try:
            size = await asyncio.to_thread(lambda: path.stat().st_size)
        except OSError as error:
            raise AttachmentError(
                f"Cannot access {_element_name(element)}."
            ) from error
        if size > max_file_bytes:
            raise AttachmentError(
                f"{_element_name(element)} exceeds the upload size limit."
            )
        payload = await asyncio.to_thread(path.read_bytes)
    if len(payload) > max_file_bytes:
        raise AttachmentError(
            f"{_element_name(element)} exceeds the upload size limit."
        )
    return payload


def _extract_text(name: str, mime: str, payload: bytes) -> str:
    extension = Path(name).suffix.casefold()
    lower_name = name.casefold()
    if extension in _ARCHIVE_EXTENSIONS or lower_name.endswith(".tar.gz"):
        return _extract_archive(name, payload)
    if extension == ".pdf" or mime == "application/pdf":
        try:
            reader = PdfReader(io.BytesIO(payload))
            return "\n\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception as error:
            raise AttachmentError(f"Could not read PDF {name}.") from error
    if (
        extension == ".docx"
        or mime
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        try:
            document = Document(io.BytesIO(payload))
            return "\n".join(paragraph.text for paragraph in document.paragraphs)
        except Exception as error:
            raise AttachmentError(f"Could not read DOCX file {name}.") from error
    if extension == ".pptx":
        try:
            presentation = Presentation(io.BytesIO(payload))
            slides = []
            for index, slide in enumerate(presentation.slides, start=1):
                text = "\n".join(
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                )
                if text:
                    slides.append(f"Slide {index}\n{text}")
            return "\n\n".join(slides)
        except Exception as error:
            raise AttachmentError(f"Could not read presentation {name}.") from error
    if extension == ".xlsx":
        try:
            workbook = load_workbook(
                io.BytesIO(payload), read_only=True, data_only=True
            )
            sheets = []
            for worksheet in workbook.worksheets:
                rows = []
                for row in worksheet.iter_rows(values_only=True):
                    values = ["" if value is None else str(value) for value in row]
                    if any(values):
                        rows.append(" | ".join(values))
                if rows:
                    sheets.append(
                        f"Sheet: {worksheet.title}\n" + "\n".join(rows)
                    )
            workbook.close()
            return "\n\n".join(sheets)
        except Exception as error:
            raise AttachmentError(f"Could not read spreadsheet {name}.") from error
    if extension in _OPEN_DOCUMENT_EXTENSIONS:
        return _extract_open_document(name, payload)
    if extension == ".epub":
        return _extract_epub(name, payload)
    if extension == ".eml" or mime == "message/rfc822":
        try:
            message = BytesParser(policy=policy.default).parsebytes(payload)
            headers = [
                f"{key}: {message.get(key)}"
                for key in ("From", "To", "Cc", "Date", "Subject")
                if message.get(key)
            ]
            bodies = []
            parts = message.walk() if message.is_multipart() else [message]
            for part in parts:
                if part.get_content_type() == "text/plain":
                    bodies.append(part.get_content())
            return "\n".join(headers + bodies)
        except Exception as error:
            raise AttachmentError(f"Could not read email {name}.") from error
    if extension == ".rtf":
        text = payload.decode("utf-8-sig", errors="replace")
        text = re.sub(r"\\'[0-9a-fA-F]{2}", " ", text)
        text = re.sub(r"\\[a-zA-Z]+-?\d* ?", " ", text)
        return re.sub(r"[{}]", "", text)

    guessed_mime = mimetypes.guess_type(name)[0] or mime
    if (
        extension not in _TEXT_EXTENSIONS
        and not guessed_mime.startswith("text/")
        and guessed_mime not in {"application/json", "application/xml"}
    ):
        raise AttachmentError(f"Unsupported attachment type: {name}.")
    text = payload.decode("utf-8-sig", errors="replace")
    if extension in {".html", ".htm"}:
        parser = _HTMLTextExtractor()
        parser.feed(text)
        return "\n".join(parser.parts)
    if extension == ".json":
        try:
            return json.dumps(json.loads(text), indent=2, ensure_ascii=False)
        except json.JSONDecodeError:
            return text
    if extension == ".csv":
        rows = csv.reader(io.StringIO(text))
        return "\n".join(" | ".join(cell for cell in row) for row in rows)
    return text


def _validate_and_encode_image(payload: bytes) -> str:
    try:
        with Image.open(io.BytesIO(payload)) as image:
            image.seek(0)
            image.load()
            if image.width * image.height > _MAX_IMAGE_PIXELS:
                raise AttachmentError("Image dimensions are too large.")
            normalized = ImageOps.exif_transpose(image)
            normalized.thumbnail(
                (_MAX_IMAGE_EDGE, _MAX_IMAGE_EDGE),
                Image.Resampling.LANCZOS,
            )
            has_alpha = normalized.mode in {"RGBA", "LA"} or (
                normalized.mode == "P" and "transparency" in normalized.info
            )
            output = io.BytesIO()
            if has_alpha:
                normalized.convert("RGBA").save(
                    output, format="PNG", optimize=True
                )
            else:
                normalized.convert("RGB").save(
                    output,
                    format="JPEG",
                    quality=88,
                    optimize=True,
                    progressive=True,
                )
    except AttachmentError:
        raise
    except Exception as error:
        raise AttachmentError("The uploaded image is invalid.") from error
    return base64.b64encode(output.getvalue()).decode("ascii")


def _extract_open_document(name: str, payload: bytes) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(payload)) as archive:
            root = ElementTree.fromstring(archive.read("content.xml"))
        parts = [text.strip() for text in root.itertext() if text.strip()]
        return "\n".join(parts)
    except Exception as error:
        raise AttachmentError(f"Could not read OpenDocument file {name}.") from error


def _extract_epub(name: str, payload: bytes) -> str:
    try:
        parts = []
        with zipfile.ZipFile(io.BytesIO(payload)) as archive:
            for member in archive.namelist():
                if not member.casefold().endswith((".html", ".htm", ".xhtml")):
                    continue
                parser = _HTMLTextExtractor()
                parser.feed(archive.read(member).decode("utf-8", errors="replace"))
                parts.extend(parser.parts)
        return "\n".join(parts)
    except Exception as error:
        raise AttachmentError(f"Could not read EPUB file {name}.") from error


def _extract_archive(name: str, payload: bytes) -> str:
    """Extract readable members from a bounded archive without writing to disk."""
    try:
        lower_name = name.casefold()
        if lower_name.endswith(".gz") and not lower_name.endswith(
            (".tar.gz", ".tgz")
        ):
            with gzip.GzipFile(fileobj=io.BytesIO(payload)) as archive:
                expanded = archive.read(_MAX_ARCHIVE_EXPANDED_BYTES + 1)
            if len(expanded) > _MAX_ARCHIVE_EXPANDED_BYTES:
                raise AttachmentError(f"Expanded archive {name} is too large.")
            inner_name = Path(name).stem or "compressed.txt"
            return _extract_text(inner_name, "", expanded)

        members: list[tuple[str, bytes]] = []
        expanded_bytes = 0
        if lower_name.endswith(".zip"):
            with zipfile.ZipFile(io.BytesIO(payload)) as archive:
                entries = [item for item in archive.infolist() if not item.is_dir()]
                if len(entries) > _MAX_ARCHIVE_MEMBERS:
                    raise AttachmentError(
                        f"Archive {name} contains too many files."
                    )
                for entry in entries:
                    if entry.flag_bits & 0x1:
                        continue
                    expanded_bytes += entry.file_size
                    if expanded_bytes > _MAX_ARCHIVE_EXPANDED_BYTES:
                        raise AttachmentError(
                            f"Expanded archive {name} is too large."
                        )
                    members.append((Path(entry.filename).name, archive.read(entry)))
        else:
            with tarfile.open(fileobj=io.BytesIO(payload), mode="r:*") as archive:
                entries = [item for item in archive.getmembers() if item.isfile()]
                if len(entries) > _MAX_ARCHIVE_MEMBERS:
                    raise AttachmentError(
                        f"Archive {name} contains too many files."
                    )
                for entry in entries:
                    expanded_bytes += entry.size
                    if expanded_bytes > _MAX_ARCHIVE_EXPANDED_BYTES:
                        raise AttachmentError(
                            f"Expanded archive {name} is too large."
                        )
                    source = archive.extractfile(entry)
                    if source is not None:
                        members.append((Path(entry.name).name, source.read()))

        sections = []
        for member_name, member_payload in members:
            if not member_name:
                continue
            member_lower = member_name.casefold()
            if (
                Path(member_lower).suffix in _ARCHIVE_EXTENSIONS
                or member_lower.endswith(".tar.gz")
            ):
                continue
            try:
                text = _extract_text(member_name, "", member_payload)
            except AttachmentError:
                continue
            if text.strip():
                sections.append(f"File: {member_name}\n{text}")
        if not sections:
            raise AttachmentError(f"No readable files were found in {name}.")
        return "\n\n".join(sections)
    except AttachmentError:
        raise
    except (gzip.BadGzipFile, tarfile.TarError, zipfile.BadZipFile, OSError) as error:
        raise AttachmentError(f"Could not read archive {name}.") from error
